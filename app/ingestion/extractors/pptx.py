from pathlib import Path

from app.ingestion.extractors.base import Extractor


class PptxExtractor(Extractor):
    supported_extensions = (".pptx",)

    def extract(self, file_path: Path) -> str:
        from pptx import Presentation

        slides: list[str] = []
        for slide in Presentation(str(file_path)).slides:
            parts = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if parts:
                slides.append("\n".join(parts))
        return "\n\n".join(slides)
